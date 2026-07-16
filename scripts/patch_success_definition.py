#!/usr/bin/env python3
"""Add Y1-minutes success definition to explorer landing + PPT slides."""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PPT = ROOT / "bundesliga_transfer_traits.pptx"
EXPLORER = ROOT / "interactive_player_explorer.html"

NAVY = RGBColor(0x1C, 0x1C, 0x1C)
ACCENT = RGBColor(0x3D, 0x5A, 0x80)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF4, 0xF5, 0xF7)
SOFT = RGBColor(0xEE, 0xF3, 0xF8)
LINE = RGBColor(0xD0, 0xD4, 0xDB)
GREEN = RGBColor(0x1B, 0x7A, 0x45)


def bg(slide, prs, color=LIGHT):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()


def txt(slide, left, top, width, height, text, size=16, bold=False, color=NAVY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Georgia"


def bullets(slide, left, top, width, height, items, size=14):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = item
        r.font.size = Pt(size)
        r.font.color.rgb = NAVY
        r.font.name = "Arial"


def slide_title(slide) -> str:
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs) or p.text
                if t.strip():
                    return t.strip()
    return ""


def update_ppt():
    prs = Presentation(str(PPT))
    if any("Y1 MINUTES" in slide_title(s) for s in prs.slides):
        print("PPT already has success-definition slides — skipping")
        return

    blank = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(blank)
    bg(s1, prs)
    txt(s1, Inches(0.6), Inches(0.45), Inches(12), Inches(0.65),
        "SUCCESS DEFINITION  ·  YEAR-1 MINUTES", size=28, bold=True, color=ACCENT)
    txt(s1, Inches(0.6), Inches(1.05), Inches(12), Inches(0.45),
        "How we operationalize “Bundesliga success” in Phase 2 with open data.", size=15, color=GRAY)
    bullets(s1, Inches(0.75), Inches(1.55), Inches(5.9), Inches(4.8), [
        "Primary success marker: Bundesliga minutes played in the player’s first BL season (Y1 minutes)",
        "Interpretation: playing-time success — did the club trust and use the inbound player?",
        "Also reported: Y1 minutes percentile (rank in cohort) and stable-trait score (profile fit)",
        "Trait & league tests ask: which prior signals line up with more Y1 minutes?",
        "Regression Δ Y1 min = estimated extra first-season minutes vs a reference league, holding controls",
    ])
    box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.95), Inches(1.55), Inches(5.7), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = SOFT
    box.line.color.rgb = LINE
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Why minutes?"
    r.font.bold = True
    r.font.size = Pt(16)
    r.font.color.rgb = ACCENT
    r.font.name = "Georgia"
    for line in [
        "",
        "Observable in FBref / StatsBomb cohorts",
        "",
        "Directly measures opportunity & role",
        "",
        "Not the same as goals, xG, or team points",
        "",
        "Honest scope: associative scouting weights, not a forecast model",
    ]:
        p = tf.add_paragraph()
        rr = p.add_run()
        rr.text = line
        rr.font.size = Pt(13)
        rr.font.color.rgb = NAVY
        rr.font.name = "Arial"

    s2 = prs.slides.add_slide(blank)
    bg(s2, prs)
    txt(s2, Inches(0.6), Inches(0.45), Inches(12), Inches(0.65),
        "IF WE HAD THE DATA  ·  TRANSFERMARKT VALUE", size=28, bold=True, color=ACCENT)
    txt(s2, Inches(0.6), Inches(1.05), Inches(12), Inches(0.45),
        "What we would use instead of (or alongside) Y1 minutes.", size=15, color=GRAY)
    bullets(s2, Inches(0.75), Inches(1.55), Inches(5.9), Inches(4.2), [
        "Top indicator in a full club model: change in market value (Δ Transfermarkt MV)",
        "Captures performance + hype + contract context in one outcome scouts care about",
        "NYRB Axis-3 used MV; this project uses open data only — Transfermarkt bans bots",
        "So we proxy success with Y1 minutes until a licensed MV feed is available",
        "Natural extension: prior traits → Δ MV with holdout regression / ML",
    ])
    box2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.95), Inches(1.55), Inches(5.7), Inches(4.2))
    box2.fill.solid()
    box2.fill.fore_color.rgb = SOFT
    box2.line.color.rgb = LINE
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "Bottom line"
    r2.font.bold = True
    r2.font.size = Pt(16)
    r2.font.color.rgb = GREEN
    r2.font.name = "Georgia"
    for line in [
        "",
        "Y1 minutes = best open-data proxy for first-year BL success",
        "",
        "Δ Transfermarkt value = preferred outcome in a pro scouting stack",
        "",
        "Same pipeline (stability → traits → outcome) — swap the Y variable",
    ]:
        p2 = tf2.add_paragraph()
        rr = p2.add_run()
        rr.text = line
        rr.font.size = Pt(13)
        rr.font.color.rgb = NAVY
        rr.font.name = "Arial"

    # Insert after slide 12 (SUCCESS INDICATORS) → index 12
    xml = prs.slides._sldIdLst  # noqa: SLF001
    ids = list(xml)
    new_ids = ids[-2:]
    for el in new_ids:
        xml.remove(el)
    xml.insert(12, new_ids[0])
    xml.insert(13, new_ids[1])

    prs.save(str(PPT))
    alias = ROOT / "bundesliga_statsbomb_deck.pptx"
    if alias.exists():
        shutil.copy(PPT, alias)
    print(f"PPT updated → {len(prs.slides)} slides")


def update_explorer():
    html = EXPLORER.read_text()

    success_css = """
  .success-def{background:#e8f4ec;border:1px solid #a8d5b8;border-radius:8px;padding:14px 16px;margin:14px 0}
  .success-def h3{font-size:13px;text-transform:uppercase;letter-spacing:.04em;color:#1b7a45;margin:0 0 8px}
  .success-def p{font-size:14px;color:#1c1c1c;line-height:1.55;margin:0}
  .success-def .alt{font-size:12px;color:#555;margin-top:10px;line-height:1.45}
"""
    if ".success-def" not in html:
        html = html.replace(
            "  .corr-note{background:",
            success_css + "  .corr-note{background:",
        )

    success_block = """
  <div class="success-def" id="success-definition">
    <h3>How we define Bundesliga success (Phase 2)</h3>
    <p>
      Our <strong>primary success marker</strong> is <strong>minutes played in the player’s first Bundesliga season</strong> (Year-1 minutes).
      We treat more Y1 minutes as <strong>playing-time success</strong> — the club gave the inbound player a real role.
      Trait tests, league comparisons, and regression <strong>Δ Y1 min</strong> all use this outcome unless noted otherwise.
    </p>
    <p class="alt">
      We also show <strong>Y1 minutes percentile</strong> (rank in cohort) and <strong>stable-trait score</strong> (profile fit on shortlisted traits).
      <strong>Transfermarkt market-value change</strong> would be the top outcome in a pro model, but we do not scrape Transfermarkt (ToS); minutes are our open-data proxy.
    </p>
  </div>
"""
    if 'id="success-definition"' not in html:
        html = html.replace(
            '  <div class="banner" id="landing-banner"></div>',
            success_block + '\n  <div class="banner" id="landing-banner"></div>',
        )

    old_method = "<li><strong>Success outcomes:</strong> Y1 minutes · Y1 minutes percentile · stable-trait score</li>"
    new_method = (
        "<li><strong>Primary success marker:</strong> <strong>Y1 minutes</strong> — first-season Bundesliga minutes played</li>"
        "<li><strong>Secondary outcomes:</strong> Y1 minutes percentile · stable-trait score</li>"
    )
    html = html.replace(old_method, new_method)

    old_banner_js = (
        "  document.getElementById('landing-banner').textContent =\n"
        "    `Phase 2 results · FBref N=${fb.cohort_n} · StatsBomb N=${SUCCESS.summary.statsbomb.cohort_n} · Trait associations with Y1 minutes & stable-trait score`;"
    )
    new_banner_js = (
        "  document.getElementById('landing-banner').textContent =\n"
        "    `Phase 2 · Success = first-season BL minutes (Y1) · FBref N=${fb.cohort_n} · StatsBomb N=${SUCCESS.summary.statsbomb.cohort_n}`;"
    )
    html = html.replace(old_banner_js, new_banner_js)

    glossary_y1 = (
        '{"term":"Y1 minutes","definition":"Bundesliga minutes played in the player\'s first BL season in our cohort."}'
    )
    glossary_y1_new = (
        '{"term":"Y1 minutes (primary success marker)","definition":"Bundesliga minutes played in the player\'s first BL season — our main Phase 2 success proxy for playing-time opportunity."}'
    )
    html = html.replace(glossary_y1, glossary_y1_new)

    if '"term":"Transfermarkt' not in html:
        html = html.replace(
            '{"term":"Stability r","definition":"Pearson r between prior-season and BL Y1 rates',
            '{"term":"Δ Transfermarkt value (not used)","definition":"Preferred pro scouting outcome — change in market value. We do not scrape Transfermarkt; Y1 minutes is the open-data substitute."},{"term":"Stability r","definition":"Pearson r between prior-season and BL Y1 rates',
        )

    EXPLORER.write_text(html)
    print("Explorer landing updated")


def main():
    update_explorer()
    update_ppt()


if __name__ == "__main__":
    main()
