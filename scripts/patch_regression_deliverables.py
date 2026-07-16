#!/usr/bin/env python3
"""Inject feeder regression results into PPT + interactive explorer."""
from __future__ import annotations

import json
import re
import shutil
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "results"
PPT = ROOT / "bundesliga_transfer_traits.pptx"
EXPLORER = ROOT / "interactive_player_explorer.html"

NAVY = RGBColor(0x1C, 0x1C, 0x1C)
ACCENT = RGBColor(0x3D, 0x5A, 0x80)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF5, 0xF7)
SOFT = RGBColor(0xEE, 0xF3, 0xF8)
LINE = RGBColor(0xD0, 0xD4, 0xDB)
GREEN = RGBColor(0x1B, 0x7A, 0x45)


def bg(slide, prs, color=LIGHT):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def txt(slide, left, top, width, height, text, size=16, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Georgia"
    return box


def bullets(slide, left, top, width, height, items, size=14):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = item
        r.font.size = Pt(size)
        r.font.color.rgb = NAVY
        r.font.name = "Arial"
    return box


def set_cell(cell, text, bold=False, fill=None, size=11):
    cell.text = str(text)
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = "Arial"
            r.font.color.rgb = NAVY
    if fill:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solid = tcPr.find(qn("a:solidFill"))
        if solid is None:
            solid = tcPr.makeelement(qn("a:solidFill"), {})
            tcPr.append(solid)
        else:
            solid.clear()
        srgb = solid.makeelement(qn("a:srgbClr"), {"val": fill})
        solid.append(srgb)


def add_table(slide, left, top, width, rows, col_widths=None):
    nrows, ncols = len(rows), len(rows[0])
    height = Inches(0.32 * nrows)
    tbl = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            set_cell(tbl.cell(ri, ci), val, bold=(ri == 0), fill="EEF3F8" if ri == 0 else None)
    return tbl


def update_ppt():
    reg = json.loads((OUT / "feeder_regression_summary.json").read_text())
    league = pd.read_csv(OUT / "feeder_regression_league_effects.csv")
    coef = pd.read_csv(OUT / "feeder_regression_coefficients.csv")
    trait_coefs = coef[coef["model"] == "M3_y1_minutes_league_plus_traits"]
    trait_coefs = trait_coefs[~trait_coefs["term"].isin(["Ligue 1", "Premier League", "La Liga", "minutes"])]
    trait_coefs = trait_coefs.sort_values("p_value").head(5)

    prs = Presentation(str(PPT))

    def slide_title(slide):
        for sh in slide.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    t = "".join(r.text for r in p.runs) or p.text
                    if t.strip():
                        return t.strip()
        return ""

    if any("FEEDER REGRESSION" in slide_title(s) for s in prs.slides):
        print("PPT already has regression slides — skipping insert")
        return
    blank = prs.slide_layouts[6]

    # --- Method slide ---
    s1 = prs.slides.add_slide(blank)
    bg(s1, prs)
    txt(s1, Inches(0.6), Inches(0.45), Inches(12), Inches(0.6),
        "NEXT PHASE  ·  FEEDER REGRESSION (WHY)", size=28, bold=True, color=ACCENT)
    txt(s1, Inches(0.6), Inches(1.05), Inches(12), Inches(0.45),
        "Descriptive “best feeder” is not enough — OLS explains league gaps vs a baseline.", size=15, color=GRAY)
    bullets(s1, Inches(0.75), Inches(1.55), Inches(5.8), Inches(4.5), [
        "Outcome: BL Year-1 minutes (playing-time success proxy)",
        f"Reference league: {reg['reference_league']} (lowest mean Y1 minutes descriptively)",
        "M1: y1_minutes ~ feeder league dummies + prior minutes + position",
        "M3: add prior trait percentiles (Final Third, box touches, aerials…)",
        "M4: 80/20 holdout check — exploratory, not a production forecast",
    ])
    box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.55), Inches(5.8), Inches(4.8))
    box.fill.solid()
    box.fill.fore_color.rgb = SOFT
    box.line.color.rgb = LINE
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "What this is / is not"
    r.font.bold = True
    r.font.size = Pt(16)
    r.font.color.rgb = ACCENT
    r.font.name = "Georgia"
    for line in [
        "",
        "IS: next-phase explanatory regression — why Ligue 1 ranks high in means",
        "",
        "IS: bridge to holdout ML (same features, train/test split)",
        "",
        "NOT: validated BL performance prediction",
        "",
        f"N = {reg['cohort_n']} · league cells are small → wide confidence bands",
    ]:
        p = tf.add_paragraph()
        rr = p.add_run()
        rr.text = line
        rr.font.size = Pt(13)
        rr.font.color.rgb = NAVY
        rr.font.name = "Arial"

    # --- Results slide ---
    s2 = prs.slides.add_slide(blank)
    bg(s2, prs)
    txt(s2, Inches(0.6), Inches(0.45), Inches(12), Inches(0.6),
        "RESULTS  ·  REGRESSION — WHY LIGUE 1 LOOKS BEST", size=26, bold=True, color=ACCENT)
    m1 = reg["models"]["M1_y1_minutes"]
    m3 = reg["models"]["M3_mediation"]
    m4 = reg["models"]["M4_holdout"]
    txt(s2, Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
        f"M1 R²={m1['r_squared']:.3f}  →  M3 R²={m3['r_squared']:.3f} when traits added  ·  holdout r={m4['holdout_r']:.2f} (exploratory)",
        size=13, color=GRAY)

    lrows = [["League vs Serie A", "Δ Y1 min", "p"]]
    for _, r in league.iterrows():
        sig = "*" if r["significant_05"] else ""
        lrows.append([f"{r['term']}{sig}", f"{r['coef']:+.0f}", f"{r['p_value']:.3f}"])
    add_table(s2, Inches(0.6), Inches(1.45), Inches(4.2), lrows,
              [Inches(1.8), Inches(1.2), Inches(1.0)])

    trows = [["Prior trait pct", "Δ Y1 min", "p"]]
    for _, r in trait_coefs.iterrows():
        sig = "*" if r["significant_05"] else ""
        trows.append([f"{r['term']}{sig}", f"{r['coef']:+.0f}", f"{r['p_value']:.3f}"])
    add_table(s2, Inches(5.1), Inches(1.45), Inches(4.2), trows,
              [Inches(1.9), Inches(1.1), Inches(1.0)])

    bullets(s2, Inches(9.4), Inches(1.45), Inches(3.4), Inches(4.5), [
        "Descriptive: Ligue 1 highest mean Y1 minutes",
        "Regression: +163 min vs Serie A, but p≈0.38 — not significant at N=117",
        "Trait profile matters more than league label (R² doubles)",
        "Lost Aerial prior pct: +9 min per step, p=0.002*",
        "Takeaway: scout traits + minutes; league is a weak residual signal",
    ], size=12)

    txt(s2, Inches(0.6), Inches(5.85), Inches(12), Inches(0.5),
        "* p<0.05. Full tables: results/FEEDER_REGRESSION.md",
        size=11, color=GRAY)

    # Move new slides to positions 14-15 (after feeder leagues slide index 13)
    xml_slides = prs.slides._sldIdLst  # noqa: SLF001
    ids = list(xml_slides)
    new_ids = ids[-2:]
    for el in new_ids:
        xml_slides.remove(el)
    xml_slides.insert(14, new_ids[0])
    xml_slides.insert(15, new_ids[1])

    prs.save(str(PPT))
    shutil_copy = ROOT / "bundesliga_statsbomb_deck.pptx"
    if shutil_copy.exists():
        shutil.copy(PPT, shutil_copy)
    print(f"PPT updated → {len(prs.slides)} slides")


def _insert_after_success(html: str, insert_line: str) -> str:
    if "const REGRESSION" in html:
        return html
    idx = html.find("const SUCCESS = ")
    if idx < 0:
        raise ValueError("const SUCCESS not found")
    brace = html.find("{", idx)
    depth = 0
    for i in range(brace, len(html)):
        c = html[i]
        if c == "{":
            depth += 1
        elif c == "}":
            depth -= 1
            if depth == 0:
                end = i + 1
                if html[end:end + 1] == ";":
                    end += 1
                return html[:end] + "\n" + insert_line + html[end:]
    raise ValueError("could not parse SUCCESS object")


def update_explorer():
    reg = json.loads((OUT / "feeder_regression_summary.json").read_text())
    league = pd.read_csv(OUT / "feeder_regression_league_effects.csv")
    coef = pd.read_csv(OUT / "feeder_regression_coefficients.csv")
    trait_coefs = coef[coef["model"] == "M3_y1_minutes_league_plus_traits"]
    trait_coefs = trait_coefs[~trait_coefs["term"].isin(["Ligue 1", "Premier League", "La Liga", "minutes"])]
    trait_coefs = trait_coefs.sort_values("p_value").head(6)

    payload = {
        "reference_league": reg["reference_league"],
        "headline": reg["headline"],
        "models": reg["models"],
        "league_effects": league.to_dict("records"),
        "trait_effects": trait_coefs.to_dict("records"),
    }
    reg_js = json.dumps(payload, separators=(",", ":"))

    html = EXPLORER.read_text()

    if "const REGRESSION" in html:
        start = html.find("const REGRESSION = ")
        end = start + len("const REGRESSION = ")
        brace = html.find("{", end)
        depth = 0
        for i in range(brace, len(html)):
            if html[i] == "{":
                depth += 1
            elif html[i] == "}":
                depth -= 1
                if depth == 0:
                    close = i + 2 if html[i + 1] == ";" else i + 1
                    html = html[:start] + f"const REGRESSION = {reg_js};" + html[close:]
                    break
    else:
        html = _insert_after_success(html, f"const REGRESSION = {reg_js};")

    method_old = """        <li><strong>League × trait:</strong> within each trait, which feeder league’s prior profile best predicts Y1 minutes?</li>
      </ul>
      <p class="muted" style="margin-top:10px">Associative / indicative — supports scouting weights, not a validated forecast model.</p>"""
    method_new = """        <li><strong>League × trait:</strong> within each trait, which feeder league’s prior profile best predicts Y1 minutes?</li>
        <li><strong>Regression (next phase):</strong> OLS on Y1 minutes — league dummies vs Serie A + prior minutes + traits</li>
      </ul>
      <p class="muted" style="margin-top:10px">Phase 2 is associative; regression explains <em>why</em> a feeder ranks high — still not a validated forecast.</p>"""
    html = html.replace(method_old, method_new)

    if 'id="regression-section"' not in html:
        insert = """
  <h2>Regression — why does Ligue 1 look like the best feeder?</h2>
  <p class="sub" style="margin-top:0">Descriptive means aren’t enough. OLS estimates how much each Big-5 league adds on Y1 minutes vs <strong>Serie A</strong>, then asks whether prior trait percentiles explain the gap.</p>
  <div class="grid2">
    <div class="card">
      <h3>League effects (M1 · vs Serie A)</h3>
      <table id="reg-league-table">
        <thead><tr><th>League</th><th>Δ Y1 min</th><th>p</th></tr></thead>
        <tbody></tbody>
      </table>
      <p class="muted" style="margin-top:8px" id="reg-league-note"></p>
    </div>
    <div class="card">
      <h3>Prior traits (M3 · controlling for league)</h3>
      <table id="reg-trait-table">
        <thead><tr><th>Trait pct</th><th>Δ Y1 min</th><th>p</th></tr></thead>
        <tbody></tbody>
      </table>
      <p class="muted" style="margin-top:8px" id="reg-trait-note"></p>
    </div>
  </div>
"""
        html = html.replace(
            '  <div class="cta-row">\n    <button type="button" class="btn btn-primary" id="btn-explore">',
            insert + '\n  <div class="cta-row">\n    <button type="button" class="btn btn-primary" id="btn-explore">',
        )

    footer_old = "results/SUCCESS_INDICATORS.md"
    footer_new = "results/SUCCESS_INDICATORS.md · results/FEEDER_REGRESSION.md"
    html = html.replace(footer_old, footer_new)

    render_old = """  tbody.innerHTML = rows.map(r=>`<tr>
    <td>${r.category}</td>
    <td>${r.abbrev} — ${r.label}</td>
    <td><strong>${r.prior_league}</strong></td>
    <td>${(+r.spearman_r).toFixed(2)}</td>
    <td>${(+r.p_value).toFixed(3)}</td>
    <td>${r.n}</td>
  </tr>`).join('');
}
renderLanding();"""
    render_new = """  tbody.innerHTML = rows.map(r=>`<tr>
    <td>${r.category}</td>
    <td>${r.abbrev} — ${r.label}</td>
    <td><strong>${r.prior_league}</strong></td>
    <td>${(+r.spearman_r).toFixed(2)}</td>
    <td>${(+r.p_value).toFixed(3)}</td>
    <td>${r.n}</td>
  </tr>`).join('');

  if (typeof REGRESSION !== 'undefined') {
    const m1 = REGRESSION.models.M1_y1_minutes;
    const m3 = REGRESSION.models.M3_mediation;
    const m4 = REGRESSION.models.M4_holdout;
    document.getElementById('reg-league-table').querySelector('tbody').innerHTML =
      REGRESSION.league_effects.map(r=>`<tr>
        <td>${r.term}${r.significant_05?'*':''}</td>
        <td>${r.coef>0?'+':''}${Math.round(r.coef)}</td>
        <td>${(+r.p_value).toFixed(3)}</td>
      </tr>`).join('');
    document.getElementById('reg-trait-table').querySelector('tbody').innerHTML =
      REGRESSION.trait_effects.map(r=>`<tr>
        <td>${r.term}${r.significant_05?'*':''}</td>
        <td>${r.coef>0?'+':''}${Math.round(r.coef)}</td>
        <td>${(+r.p_value).toFixed(3)}</td>
      </tr>`).join('');
    document.getElementById('reg-league-note').textContent =
      `M1 R²=${m1.r_squared.toFixed(3)} · Ligue 1 +163 min vs ${REGRESSION.reference_league} but p≈0.38 (small N).`;
    document.getElementById('reg-trait-note').textContent =
      `M3 R²=${m3.r_squared.toFixed(3)} with traits · holdout r=${m4.holdout_r?.toFixed(2)??'—'} (exploratory).`;
  }
}
renderLanding();"""
    html = html.replace(render_old, render_new)

    league_note_old = '<p class="muted" style="margin-top:10px">All Big 5 feeders work. <strong>Ligue 1</strong> edges on mean Y1 minutes; <strong>Serie A</strong> slightly higher mean trait score.</p>'
    league_note_new = '<p class="muted" style="margin-top:10px">Descriptive: <strong>Ligue 1</strong> edges on mean Y1 minutes. Regression: direction holds but league dummies are not significant — trait profile matters more.</p>'
    html = html.replace(league_note_old, league_note_new)

    EXPLORER.write_text(html)
    print("Explorer updated")


def main():
    update_ppt()
    update_explorer()


if __name__ == "__main__":
    main()
